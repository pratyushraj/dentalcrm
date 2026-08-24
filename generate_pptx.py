import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    NAVY = RGBColor(11, 36, 80)       # #0B2450
    BLUE = RGBColor(8, 103, 232)      # #0867E8
    TEAL = RGBColor(18, 168, 160)     # #12A8A0
    DARK_BG = RGBColor(15, 23, 42)    # #0F172A
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252) # #F8FAFC
    GRAY_TEXT = RGBColor(100, 116, 139) # #64748B
    BORDER_COLOR = RGBColor(226, 232, 240) # #E2E8F0
    ACCENT_BG = RGBColor(238, 242, 255)

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    def add_header(slide, tag_text, title_text):
        # Category Tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # Slide Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = NAVY

    # SLIDE 1: COVER
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1, LIGHT_BG)

    # Decorative Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER_COLOR

    # Badge
    tx = slide1.shapes.add_textbox(Inches(1.3), Inches(1.3), Inches(5), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "PARTNERSHIP OPPORTUNITY WITH VROZART"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEAL

    # Main Title
    tx = slide1.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(9), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "Clinaza"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tx.text_frame.add_paragraph()
    p2.text = '"EMI for Better Health"'
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = BLUE

    # Subtitle
    tx = slide1.shapes.add_textbox(Inches(1.3), Inches(3.4), Inches(9), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = "Healthcare financing infrastructure for clinics and patients"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_TEXT

    # Footer
    tx = slide1.shapes.add_textbox(Inches(1.3), Inches(5.8), Inches(10), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "Strategic Deck for CEO & Senior Leadership Team · Vrozart × Clinaza"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT

    # SLIDE 2: THE PROBLEM
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2, WHITE)
    add_header(slide2, "01 / Market Friction", "The Healthcare Affordability Gap")

    probs = [
        ("01", "Patients Defer Care", "Planned dental procedures (implants, aligners, crowns) require ₹30,000 to ₹3,00,000 upfront, forcing patients to delay or reject treatment."),
        ("02", "Clinics Lose Conversions", "Doctors spend consultation time diagnosing, only to lose treatment acceptance at the front desk when patients cannot pay lump-sum."),
        ("03", "Lenders Lack Distribution", "Banks & NBFCs have capital but lack point-of-care digital origination at the doctor's desk to acquire high-intent healthcare credit demand.")
    ]

    for i, (num, title, desc) in enumerate(probs):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        w = Inches(3.7)
        h = Inches(4.5)
        
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = BORDER_COLOR

        tx = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.4), w - Inches(0.6), h - Inches(0.8))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(14)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY_TEXT
        p3.space_before = Pt(10)

    # SLIDE 3: SOLUTION
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3, WHITE)
    add_header(slide3, "02 / Our Architecture", "The Clinaza Operating Model")

    steps = [
        ("STEP 01", "Doctor Consultation", "Treatment Plan Created"),
        ("STEP 02", "Clinaza Pre-Check", "2-Min Digital Assessment"),
        ("STEP 03", "Lender API Underwriting", "Sanction & Agreement"),
        ("STEP 04", "Disbursement & EMI", "Treatment Starts Immediately")
    ]

    for i, (s_num, s_title, s_sub) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        top = Inches(1.8)
        w = Inches(2.8)
        h = Inches(1.8)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = BORDER_COLOR

        tx = slide3.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), w - Inches(0.3), h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_num
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p2 = tf.add_paragraph()
        p2.text = s_title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        
        p3 = tf.add_paragraph()
        p3.text = s_sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRAY_TEXT

    # Dark Banner on Slide 3
    banner = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_BG
    banner.line.fill.background()

    tx = slide3.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11.0), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STRICT OPERATIONAL BOUNDARIES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p2 = tf.add_paragraph()
    p2.text = "Clinaza does NOT become the lender. We do not take credit risk or hold a balance sheet. Regulated lending partners (Banks/NBFCs) handle credit policy, underwriting, sanction, loan agreement, and servicing. Clinaza focuses purely on clinic distribution, software infrastructure, and patient financing demand."
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)

    # SLIDE 4: TRACTION
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4, WHITE)
    add_header(slide4, "03 / Early Validation", "Current Network Traction")

    metrics = [
        ("20+", "Clinics Onboarded", "Active dental network"),
        ("Expanding", "Clinic Network", "Active doctor onboarding"),
        ("Incoming", "Patient Loan Requests", "Active pre-check files"),
        ("₹30K–₹3L", "Ticket Requirement", "Planned dental procedures")
    ]

    for i, (val, title, sub) in enumerate(metrics):
        left = Inches(0.8 + i * 2.95)
        top = Inches(1.8)
        w = Inches(2.8)
        h = Inches(3.2)

        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BG
        box.line.color.rgb = BLUE

        tx = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), w - Inches(0.4), h - Inches(1.0))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = val
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(14)

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRAY_TEXT

    # Focus Footer
    tx = slide4.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "INITIAL FOCUS VERTICAL: DENTAL CARE PROCEDURES"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Primary use cases: Implants, clear aligners, ceramic braces, crowns, and full mouth rehabilitations."
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY_TEXT

    # SLIDE 5: WHY THIS CAN SCALE
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5, WHITE)
    add_header(slide5, "04 / Flywheel", "The Healthcare Distribution Flywheel")

    # Flywheel Box
    fbox = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    fbox.fill.solid()
    fbox.fill.fore_color.rgb = LIGHT_BG
    fbox.line.color.rgb = BORDER_COLOR

    tx = slide5.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "More Onboarded Clinics  ➔  More Patient EMI Requests  ➔  More Lender Approvals  ➔  Better Conversion  ➔  Scale Network"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # 2 Sub Cards
    sub_cards = [
        ("Point-of-Care Advantage", "Financing intent is captured directly inside the consultation room when the treatment decision is made — yielding higher conversion than generic digital personal loans."),
        ("Scalable Distribution Layer", "Clinaza standardizes clinic onboarding and digital pre-checks, allowing regulated lenders to access thousands of doctors through a single infrastructure integration.")
    ]

    for i, (ctitle, cdesc) in enumerate(sub_cards):
        left = Inches(0.8 + i * 5.95)
        top = Inches(4.3)
        w = Inches(5.75)
        h = Inches(2.4)

        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER_COLOR

        tx = slide5.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), w - Inches(0.6), h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_before = Pt(8)

    # SLIDE 6: WHAT CLINAZA NEEDS FROM VROZART
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6, DARK_BG)

    # Header for dark slide
    txBox = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "05 / STRATEGIC ALIGNMENT".upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = TEAL

    txBox2 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "What We Are Seeking from Vrozart"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    needs = [
        ("Lending Partner Access", "Introductions to Banks, NBFCs & capital partners for ₹30K–₹3L tickets."),
        ("Suitable Credit Products", "Structuring personal loan / healthcare EMI products fitting patient profiles."),
        ("API & Tech Infrastructure", "Co-building or leveraging existing embedded finance APIs."),
        ("Underwriting Guidance", "Expertise on credit policy, bureau checks & risk boundaries."),
        ("Compliance & Workflow", "Guidance on digital lending compliance (LSP / FLDG / RBI rules)."),
        ("Strategic Mentorship", "Founder guidance on scaling a regulated lending ecosystem.")
    ]

    for i, (ntitle, ndesc) in enumerate(needs):
        row = i // 3
        col = i % 3
        left = Inches(0.8 + col * 3.9)
        top = Inches(1.8 + row * 2.3)
        w = Inches(3.7)
        h = Inches(2.0)

        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = RGBColor(51, 65, 85)

        tx = slide6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), w - Inches(0.4), h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ntitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEAL

        p2 = tf.add_paragraph()
        p2.text = ndesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(203, 213, 225)
        p2.space_before = Pt(6)

    # SLIDE 7: TECH CAPABILITY
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7, WHITE)
    add_header(slide7, "06 / Tech Capability", "Developer-Led Integration Architecture")

    # Flow Box
    fbox = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0))
    fbox.fill.solid()
    fbox.fill.fore_color.rgb = LIGHT_BG
    fbox.line.color.rgb = BORDER_COLOR

    tx = slide7.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Clinic / Patient Application  ➔  KYC & Consent  ➔  Lender API Engine  ➔  Bureau & Sanction  ➔  eSign & Disbursement"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Developer Note
    dev_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.5))
    dev_box.fill.solid()
    dev_box.fill.fore_color.rgb = ACCENT_BG
    dev_box.line.color.rgb = BLUE

    tx = slide7.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(11.1), Inches(2.1))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DEVELOPER-LED INTEGRATION AGILITY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    p2 = tf.add_paragraph()
    p2.text = "Clinaza is developer-led. We can seamlessly integrate via REST APIs, webhooks, or SDKs to connect clinic front-desks directly into a lender's underwriting engine. We support real-time status sync back to doctors and patients."
    p2.font.size = Pt(13)
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(8)

    # SLIDE 8: PARTNERSHIP MODEL
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8, WHITE)
    add_header(slide8, "07 / Operational Split", "Clear Responsibility Division")

    roles = [
        ("CLINAZA ROLE", [
            "• 20+ clinic network growth",
            "• Clinic onboarding & support",
            "• Patient application journey",
            "• Technical API integration"
        ], BLUE),
        ("LENDING PARTNER", [
            "• Credit policy & underwriting",
            "• Bureau check & loan approval",
            "• Sanction & loan agreement",
            "• Disbursement & EMI servicing"
        ], TEAL),
        ("VROZART PARTNER", [
            "• Lender introductions",
            "• Product structuring guidance",
            "• Finance ecosystem access",
            "• Strategic scale partnership"
        ], NAVY)
    ]

    for i, (rtitle, items, color) in enumerate(roles):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        w = Inches(3.7)
        h = Inches(4.8)

        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = BORDER_COLOR

        tx = slide8.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), w - Inches(0.6), h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        for item in items:
            p2 = tf.add_paragraph()
            p2.text = item
            p2.font.size = Pt(12)
            p2.font.color.rgb = GRAY_TEXT
            p2.space_before = Pt(10)

    # SLIDE 9: COMMERCIAL MODEL
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9, WHITE)
    add_header(slide9, "08 / Monetization Structure", "Commercial Models & Principles")

    left_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_BG
    left_box.line.color.rgb = BORDER_COLOR

    tx = slide9.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "COMMERCIAL OPTIONS FOR DISCUSSION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    opts = [
        ("01. Referral / Origination Fee", "LSP origination fee per sanctioned or disbursed loan file."),
        ("02. Disbursement Revenue Share", "Percentage share on total monthly disbursed volume."),
        ("03. Strategic Joint Partnership", "Co-building healthcare financing distribution with shared economics.")
    ]

    for o_title, o_desc in opts:
        p2 = tf.add_paragraph()
        p2.text = o_title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = o_desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRAY_TEXT

    # Right Principles Card (Dark)
    right_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = DARK_BG
    right_box.line.fill.background()

    tx = slide9.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PATIENT & CLINIC PRINCIPLES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p2 = tf.add_paragraph()
    p2.text = "₹0 Fee to Dental Clinics"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "Clinaza does NOT charge clinics any financing setup or subscription fee. The lender charges the borrower according to approved loan terms.\n\nCommercial arrangements between Clinaza and Vrozart / lending partners can be finalized upon mutual discussion."
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(10)

    # SLIDE 10: THE ASK & NEXT STEPS
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10, LIGHT_BG)
    add_header(slide10, "09 / Next Steps", "Let's Build the Healthcare Financing Layer Together")

    # Left Steps
    lbox = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = WHITE
    lbox.line.color.rgb = BORDER_COLOR

    tx = slide10.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PROPOSED ACTION PLAN"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    actions = [
        "1. Identify 1–2 suitable lending partners / NBFC channels",
        "2. Evaluate Clinaza's 20+ clinic network pipeline & demand",
        "3. Define product eligibility for ₹30K–₹3L ticket sizes",
        "4. Agree on a pilot starting with existing 20+ clinics",
        "5. Scale conversion & disbursement across network"
    ]
    for act in actions:
        p2 = tf.add_paragraph()
        p2.text = act
        p2.font.size = Pt(12)
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(10)

    # Right Closing Card (Navy)
    rbox = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    rbox.fill.solid()
    rbox.fill.fore_color.rgb = NAVY
    rbox.line.fill.background()

    tx = rbox.text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = "FOUNDER CLOSING"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p2 = tx.add_paragraph()
    p2.text = '"Clinaza brings the healthcare distribution. We are looking for the right financial and strategic partner in Vrozart to build the lending layer."'
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(14)

    p3 = tx.add_paragraph()
    p3.text = "\nContact & Web:\n• Website: clinaza.in\n• WhatsApp: 7292984244"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(191, 219, 254)
    p3.space_before = Pt(14)

    output_path = "/Users/pratyushraj/Desktop/dental-crm/public/Clinaza_Vrozart_Partnership_Deck.pptx"
    prs.save(output_path)
    print(f"PPTX successfully created at: {output_path}")

if __name__ == "__main__":
    create_deck()
